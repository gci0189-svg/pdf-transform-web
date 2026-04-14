"""
DocFlow 文件轉換後端 — FastAPI
Railway 部署版：自動讀取 $PORT、提供靜態前端
"""

import os
import uuid
import shutil
import zipfile
import subprocess
from pathlib import Path
from typing import List

from fastapi import FastAPI, File, UploadFile, Form, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, HTMLResponse
from fastapi.staticfiles import StaticFiles

# ── 初始化 ───────────────────────────────────────────────────
app = FastAPI(title="DocFlow API", version="1.0.0")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)

BASE_DIR   = Path(__file__).parent
UPLOAD_DIR = BASE_DIR / "uploads"
OUTPUT_DIR = BASE_DIR / "outputs"
STATIC_DIR = BASE_DIR / "static"

for d in (UPLOAD_DIR, OUTPUT_DIR, STATIC_DIR):
    d.mkdir(exist_ok=True)

# index.html 複製到 static/
_src = BASE_DIR / "index.html"
_dst = STATIC_DIR / "index.html"
if _src.exists() and not _dst.exists():
    shutil.copy(_src, _dst)

# 掛載靜態目錄
app.mount("/static", StaticFiles(directory=str(STATIC_DIR)), name="static")


# ── 工具函式 ─────────────────────────────────────────────────
def save_upload(file: UploadFile) -> Path:
    suffix = Path(file.filename).suffix
    dest = UPLOAD_DIR / f"{uuid.uuid4()}{suffix}"
    with dest.open("wb") as f:
        shutil.copyfileobj(file.file, f)
    return dest

def make_output_path(suffix: str) -> Path:
    return OUTPUT_DIR / f"{uuid.uuid4()}{suffix}"

def cleanup(*paths: Path):
    for p in paths:
        try:
            if p and p.exists():
                p.unlink()
        except Exception:
            pass


# ── 根路徑 → 前端 HTML ───────────────────────────────────────
@app.get("/", response_class=HTMLResponse)
async def root():
    # 優先讀取 static/index.html，fallback 到根目錄
    for candidate in (STATIC_DIR / "index.html", BASE_DIR / "index.html"):
        if candidate.exists():
            return HTMLResponse(content=candidate.read_text(encoding="utf-8"))
    return HTMLResponse(content="<h1>DocFlow</h1><p>找不到 index.html</p>", status_code=404)


# ── 健康檢查 ─────────────────────────────────────────────────
@app.get("/api/health")
def health():
    return {"status": "ok", "message": "DocFlow API 運行中"}


# ── 1. PDF → Word ────────────────────────────────────────────
@app.post("/api/convert/pdf-to-word")
async def pdf_to_word(
    file: UploadFile = File(...),
    ocr: str = Form("false"),
    keep_images: str = Form("true"),
):
    if not file.filename.lower().endswith(".pdf"):
        raise HTTPException(400, "請上傳 PDF 檔案")
    src = save_upload(file)
    out = make_output_path(".docx")
    try:
        from pdf2docx import Converter
        cv = Converter(str(src))
        cv.convert(str(out), start=0, end=None)
        cv.close()
        return FileResponse(
            path=str(out),
            filename=Path(file.filename).stem + ".docx",
            media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        )
    except ImportError:
        raise HTTPException(500, "缺少 pdf2docx")
    except Exception as e:
        cleanup(out)
        raise HTTPException(500, f"轉換失敗：{e}")
    finally:
        cleanup(src)


# ── 2. Word → PDF ────────────────────────────────────────────
@app.post("/api/convert/word-to-pdf")
async def word_to_pdf(file: UploadFile = File(...)):
    if not file.filename.lower().endswith((".docx", ".doc")):
        raise HTTPException(400, "請上傳 Word 檔案")
    src = save_upload(file)
    out_dir = OUTPUT_DIR / uuid.uuid4().hex
    out_dir.mkdir()
    try:
        result = subprocess.run(
            ["libreoffice", "--headless", "--convert-to", "pdf",
             "--outdir", str(out_dir), str(src)],
            capture_output=True, text=True, timeout=60
        )
        if result.returncode != 0:
            raise Exception(result.stderr)
        pdf_files = list(out_dir.glob("*.pdf"))
        if not pdf_files:
            raise Exception("LibreOffice 未產生 PDF")
        return FileResponse(
            path=str(pdf_files[0]),
            filename=Path(file.filename).stem + ".pdf",
            media_type="application/pdf",
        )
    except FileNotFoundError:
        raise HTTPException(500, "LibreOffice 未安裝")
    except Exception as e:
        raise HTTPException(500, f"轉換失敗：{e}")
    finally:
        cleanup(src)


# ── 3. PDF → PPT ─────────────────────────────────────────────
@app.post("/api/convert/pdf-to-ppt")
async def pdf_to_ppt(
    file: UploadFile = File(...),
    dpi: str = Form("150"),
    aspect: str = Form("16:9"),
):
    if not file.filename.lower().endswith(".pdf"):
        raise HTTPException(400, "請上傳 PDF 檔案")
    src = save_upload(file)
    out = make_output_path(".pptx")
    try:
        from pdf2image import convert_from_path
        from pptx import Presentation
        from pptx.util import Inches
        import io
        images = convert_from_path(str(src), dpi=int(dpi))
        prs = Presentation()
        prs.slide_width  = Inches(13.33 if aspect == "16:9" else 10)
        prs.slide_height = Inches(7.5)
        blank = prs.slide_layouts[6]
        for img in images:
            slide = prs.slides.add_slide(blank)
            buf = io.BytesIO()
            img.save(buf, format="PNG")
            buf.seek(0)
            slide.shapes.add_picture(buf, 0, 0, width=prs.slide_width, height=prs.slide_height)
        prs.save(str(out))
        return FileResponse(
            path=str(out),
            filename=Path(file.filename).stem + ".pptx",
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
    except Exception as e:
        cleanup(out)
        raise HTTPException(500, f"轉換失敗：{e}")
    finally:
        cleanup(src)


# ── 4. PDF → Excel ───────────────────────────────────────────
@app.post("/api/convert/pdf-to-excel")
async def pdf_to_excel(file: UploadFile = File(...)):
    if not file.filename.lower().endswith(".pdf"):
        raise HTTPException(400, "請上傳 PDF 檔案")
    src = save_upload(file)
    out = make_output_path(".xlsx")
    try:
        import pdfplumber, openpyxl
        wb = openpyxl.Workbook()
        wb.remove(wb.active)
        with pdfplumber.open(str(src)) as pdf:
            for pn, page in enumerate(pdf.pages, 1):
                tables = page.extract_tables()
                for ti, table in enumerate(tables):
                    if not table:
                        continue
                    ws = wb.create_sheet(title=f"P{pn}_T{ti+1}")
                    for row in table:
                        ws.append([cell or "" for cell in row])
                if not tables:
                    text = page.extract_text() or ""
                    if text.strip():
                        ws = wb.create_sheet(title=f"P{pn}_Text")
                        for line in text.splitlines():
                            ws.append([line])
        if not wb.sheetnames:
            wb.create_sheet("Empty")
        wb.save(str(out))
        return FileResponse(
            path=str(out),
            filename=Path(file.filename).stem + ".xlsx",
            media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        )
    except Exception as e:
        cleanup(out)
        raise HTTPException(500, f"轉換失敗：{e}")
    finally:
        cleanup(src)


# ── 5. 圖片格式轉換 ──────────────────────────────────────────
@app.post("/api/convert/image")
async def image_convert(
    file: UploadFile = File(...),
    output_format: str = Form("png"),
    quality: str = Form("90"),
):
    src = save_upload(file)
    out_ext = "jpg" if output_format == "jpeg" else output_format
    out = make_output_path(f".{out_ext}")
    try:
        from PIL import Image
        img = Image.open(str(src))
        if output_format in ("jpg", "jpeg") and img.mode in ("RGBA", "P"):
            bg = Image.new("RGB", img.size, (255, 255, 255))
            if img.mode == "P":
                img = img.convert("RGBA")
            bg.paste(img, mask=img.split()[3] if img.mode == "RGBA" else None)
            img = bg
        elif output_format in ("png", "webp") and img.mode == "P":
            img = img.convert("RGBA")
        save_kwargs = {}
        if output_format in ("jpg", "jpeg", "webp"):
            save_kwargs["quality"] = max(1, min(95, int(quality)))
        img.save(str(out), format=output_format.upper(), **save_kwargs)
        mime = {"jpg":"image/jpeg","jpeg":"image/jpeg","png":"image/png",
                "webp":"image/webp","bmp":"image/bmp","tiff":"image/tiff"}
        return FileResponse(
            path=str(out),
            filename=Path(file.filename).stem + f".{out_ext}",
            media_type=mime.get(output_format, "application/octet-stream"),
        )
    except Exception as e:
        cleanup(out)
        raise HTTPException(500, f"轉換失敗：{e}")
    finally:
        cleanup(src)


# ── 6. PDF 壓縮 ──────────────────────────────────────────────
@app.post("/api/pdf/compress")
async def pdf_compress(
    file: UploadFile = File(...),
    level: str = Form("medium"),
):
    if not file.filename.lower().endswith(".pdf"):
        raise HTTPException(400, "請上傳 PDF 檔案")
    src = save_upload(file)
    out = make_output_path(".pdf")
    gs_setting = {"low":"/printer","medium":"/ebook","high":"/screen"}.get(level, "/ebook")
    try:
        result = subprocess.run([
            "gs", "-sDEVICE=pdfwrite", "-dCompatibilityLevel=1.4",
            f"-dPDFSETTINGS={gs_setting}", "-dNOPAUSE", "-dQUIET", "-dBATCH",
            f"-sOutputFile={out}", str(src)
        ], capture_output=True, text=True, timeout=120)
        if result.returncode != 0:
            raise Exception(result.stderr)
        original   = src.stat().st_size
        compressed = out.stat().st_size
        ratio = round((1 - compressed / original) * 100, 1)
        resp = FileResponse(
            path=str(out),
            filename=Path(file.filename).stem + "_compressed.pdf",
            media_type="application/pdf",
        )
        resp.headers["X-Original-Size"]    = str(original)
        resp.headers["X-Compressed-Size"]  = str(compressed)
        resp.headers["X-Compression-Ratio"] = str(ratio)
        resp.headers["Access-Control-Expose-Headers"] = (
            "X-Original-Size, X-Compressed-Size, X-Compression-Ratio"
        )
        return resp
    except FileNotFoundError:
        raise HTTPException(500, "Ghostscript 未安裝")
    except Exception as e:
        cleanup(out)
        raise HTTPException(500, f"壓縮失敗：{e}")
    finally:
        cleanup(src)


# ── 7. PDF 合併 ──────────────────────────────────────────────
@app.post("/api/pdf/merge")
async def pdf_merge(files: List[UploadFile] = File(...)):
    if len(files) < 2:
        raise HTTPException(400, "請至少上傳兩個 PDF 檔案")
    srcs, out = [], make_output_path(".pdf")
    try:
        from pypdf import PdfWriter, PdfReader
        writer = PdfWriter()
        for f in files:
            if not f.filename.lower().endswith(".pdf"):
                raise HTTPException(400, f"{f.filename} 不是 PDF")
            src = save_upload(f)
            srcs.append(src)
            for page in PdfReader(str(src)).pages:
                writer.add_page(page)
        with open(str(out), "wb") as fp:
            writer.write(fp)
        return FileResponse(path=str(out), filename="merged.pdf", media_type="application/pdf")
    except Exception as e:
        cleanup(out)
        raise HTTPException(500, f"合併失敗：{e}")
    finally:
        for s in srcs:
            cleanup(s)


# ── 8. PDF 分割 ──────────────────────────────────────────────
@app.post("/api/pdf/split")
async def pdf_split(
    file: UploadFile = File(...),
    mode: str = Form("each"),
    page_range: str = Form("1-3"),
):
    if not file.filename.lower().endswith(".pdf"):
        raise HTTPException(400, "請上傳 PDF 檔案")
    src = save_upload(file)
    zip_out = make_output_path(".zip")
    tmp_dir = OUTPUT_DIR / uuid.uuid4().hex
    tmp_dir.mkdir()
    try:
        from pypdf import PdfWriter, PdfReader
        reader = PdfReader(str(src))
        total  = len(reader.pages)

        def parse_range(s, total):
            pages = set()
            for part in s.split(","):
                part = part.strip()
                if "-" in part:
                    a, b = part.split("-", 1)
                    pages.update(range(int(a)-1, min(int(b), total)))
                elif part.isdigit():
                    p = int(part) - 1
                    if 0 <= p < total:
                        pages.add(p)
            return sorted(pages)

        indices = list(range(total)) if mode == "each" else parse_range(page_range, total)
        for idx in indices:
            w = PdfWriter()
            w.add_page(reader.pages[idx])
            with open(str(tmp_dir / f"page_{idx+1:03d}.pdf"), "wb") as fp:
                w.write(fp)
        with zipfile.ZipFile(str(zip_out), "w", zipfile.ZIP_DEFLATED) as zf:
            for pdf in sorted(tmp_dir.glob("*.pdf")):
                zf.write(pdf, pdf.name)
        return FileResponse(
            path=str(zip_out),
            filename=Path(file.filename).stem + "_split.zip",
            media_type="application/zip",
        )
    except Exception as e:
        raise HTTPException(500, f"分割失敗：{e}")
    finally:
        cleanup(src)
        shutil.rmtree(tmp_dir, ignore_errors=True)


# ── 9. OCR ──────────────────────────────────────────────────
@app.post("/api/convert/ocr")
async def ocr_convert(
    file: UploadFile = File(...),
    lang: str = Form("chi_tra+eng"),
    output_format: str = Form("txt"),
):
    src = save_upload(file)
    suffix = {"txt":".txt","docx":".docx","pdf":".pdf"}.get(output_format, ".txt")
    out = make_output_path(suffix)
    try:
        import pytesseract
        from PIL import Image
        ext = Path(file.filename).suffix.lower()
        if ext == ".pdf":
            from pdf2image import convert_from_path
            images = convert_from_path(str(src), dpi=300)
        else:
            images = [Image.open(str(src))]
        full_text = "\n\n".join(pytesseract.image_to_string(img, lang=lang) for img in images)
        if output_format == "txt":
            out.write_text(full_text, encoding="utf-8")
        elif output_format == "docx":
            from docx import Document
            doc = Document()
            for para in full_text.split("\n"):
                doc.add_paragraph(para)
            doc.save(str(out))
        else:
            from reportlab.lib.pagesizes import A4
            from reportlab.pdfgen import canvas as rl_canvas
            c = rl_canvas.Canvas(str(out), pagesize=A4)
            w, h = A4
            y = h - 40
            for line in full_text.splitlines():
                if y < 40:
                    c.showPage(); y = h - 40
                c.drawString(40, y, line[:120])
                y -= 14
            c.save()
        mime = {".txt":"text/plain",
                ".docx":"application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                ".pdf":"application/pdf"}
        return FileResponse(
            path=str(out),
            filename=Path(file.filename).stem + f"_ocr{suffix}",
            media_type=mime.get(suffix, "application/octet-stream"),
        )
    except Exception as e:
        cleanup(out)
        raise HTTPException(500, f"OCR 失敗：{e}")
    finally:
        cleanup(src)


# ── 主程式 ───────────────────────────────────────────────────
if __name__ == "__main__":
    import uvicorn
    port = int(os.environ.get("PORT", 8000))
    uvicorn.run("main:app", host="0.0.0.0", port=port, reload=False)
